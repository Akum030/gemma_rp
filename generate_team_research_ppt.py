from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUT = ROOT / "priority_attribute_extraction_team_presentation_may12_2026.pptx"

BG = RGBColor(248, 247, 243)
INK = RGBColor(26, 32, 44)
ACCENT = RGBColor(12, 94, 145)
ACCENT_2 = RGBColor(178, 69, 52)
MUTED = RGBColor(94, 106, 120)


def add_bg(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def title_style(paragraph, size=28, color=INK):
    paragraph.font.size = Pt(size)
    paragraph.font.bold = True
    paragraph.font.color.rgb = color


def body_style(paragraph, size=18, color=INK):
    paragraph.font.size = Pt(size)
    paragraph.font.color.rgb = color


def configure_text_frame(text_frame, margin=0.06, word_wrap=True):
    text_frame.word_wrap = word_wrap
    text_frame.margin_left = Inches(margin)
    text_frame.margin_right = Inches(margin)
    text_frame.margin_top = Inches(margin)
    text_frame.margin_bottom = Inches(margin)


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.15))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()

    tx = slide.shapes.add_textbox(Inches(0.8), Inches(1.45), Inches(11.8), Inches(1.6))
    configure_text_frame(tx.text_frame)
    p = tx.text_frame.paragraphs[0]
    p.text = title
    title_style(p, 27)

    sub = slide.shapes.add_textbox(Inches(0.82), Inches(2.85), Inches(11.7), Inches(2.15))
    configure_text_frame(sub.text_frame)
    p2 = sub.text_frame.paragraphs[0]
    p2.text = subtitle
    body_style(p2, 18, MUTED)

    footer = slide.shapes.add_textbox(Inches(0.82), Inches(6.55), Inches(6.0), Inches(0.4))
    configure_text_frame(footer.text_frame, margin=0.02)
    p3 = footer.text_frame.paragraphs[0]
    p3.text = "Internal team deck | Verified through V11 | 11 May 2026"
    body_style(p3, 12, MUTED)
    return slide


def add_bullets_slide(prs, title, bullets, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    head = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12.0), Inches(0.7))
    configure_text_frame(head.text_frame, margin=0.02)
    p = head.text_frame.paragraphs[0]
    p.text = title
    title_style(p, 24)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(1.2), Inches(2.4), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.95), Inches(1.55), Inches(11.2), Inches(4.9))
    tf = box.text_frame
    configure_text_frame(tf)
    first = True
    for bullet in bullets:
        paragraph = tf.paragraphs[0] if first else tf.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        body_style(paragraph, 19)
        first = False

    if note:
        note_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(6.3), Inches(11.1), Inches(0.7))
        note_box.fill.solid()
        note_box.fill.fore_color.rgb = RGBColor(232, 239, 245)
        note_box.line.color.rgb = ACCENT
        configure_text_frame(note_box.text_frame, margin=0.04)
        txt = note_box.text_frame.paragraphs[0]
        txt.text = note
        body_style(txt, 14, MUTED)
    return slide


def add_table_slide(prs, title, columns, rows, note=None, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    head = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12.0), Inches(0.7))
    configure_text_frame(head.text_frame, margin=0.02)
    p = head.text_frame.paragraphs[0]
    p.text = title
    title_style(p, 24)

    table = slide.shapes.add_table(len(rows) + 1, len(columns), Inches(0.6), Inches(1.35), Inches(12.1), Inches(4.7)).table
    if col_widths and len(columns) == len(col_widths):
        for idx, width in enumerate(col_widths):
            table.columns[idx].width = Inches(width)

    for idx, label in enumerate(columns):
        cell = table.cell(0, idx)
        cell.text = label
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        configure_text_frame(cell.text_frame, margin=0.03)
        para = cell.text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        body_style(para, 13, RGBColor(255, 255, 255))
        para.font.bold = True

    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(241, 244, 247)
            configure_text_frame(cell.text_frame, margin=0.03)
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.CENTER
            body_style(para, 13)

    if note:
        box = slide.shapes.add_textbox(Inches(0.8), Inches(6.35), Inches(11.5), Inches(0.5))
        configure_text_frame(box.text_frame, margin=0.02)
        p = box.text_frame.paragraphs[0]
        p.text = note
        body_style(p, 13, MUTED)
    return slide


def add_example_slide(prs, title, left_title, left_body, right_title, right_body):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    head = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12.0), Inches(0.7))
    configure_text_frame(head.text_frame, margin=0.02)
    p = head.text_frame.paragraphs[0]
    p.text = title
    title_style(p, 24)

    for x, section_title, body, accent_color in [
        (0.7, left_title, left_body, ACCENT),
        (6.8, right_title, right_body, ACCENT_2),
    ]:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.4), Inches(5.7), Inches(5.4))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(255, 255, 255)
        card.line.color.rgb = accent_color
        title_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(1.62), Inches(5.2), Inches(0.35))
        configure_text_frame(title_box.text_frame, margin=0.02)
        p1 = title_box.text_frame.paragraphs[0]
        p1.text = section_title
        title_style(p1, 15, accent_color)
        body_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.02), Inches(5.15), Inches(4.45))
        configure_text_frame(body_box.text_frame, margin=0.05)
        p2 = body_box.text_frame.paragraphs[0]
        p2.text = body
        body_style(p2, 12)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Priority-Ordered Attribute Extraction: Complete Research Story",
        "Internal team deck covering the audited Gemma 4 program from V3b through V11, the honest Gemma 3 position, the strongest results, and what remains before publication.",
    )

    add_bullets_slide(
        prs,
        "Why This Work Matters",
        [
            "Industrial queries are not just unordered attribute bags; they contain priority and intent structure.",
            "Downstream search systems benefit from knowing which attribute should drive retrieval first and which constraints are secondary.",
            "The project formalized this as a nested JSON prediction task with attribute order and key-synonym order.",
            "The preserved evidence now covers the whole path from early noisy Gemma 4 extraction to the current V11 standalone result and the hybrid deployment path.",
        ],
        note="Core claim: priority-aware structured extraction is useful and measurably harder than ordinary flat extraction.",
    )

    add_bullets_slide(
        prs,
        "Gemma 3 And Gemma 4",
        [
            "Gemma 3 survives in the repo only as exploratory setup context via Ollama pull commands.",
            "There are no archived Gemma 3 fine-tuning logs or scored benchmark summaries, so a metric table would not be reproducible today.",
            "Gemma 4 E4B is the fully audited platform for all preserved fine-tuning scripts, evaluations, and benchmark reports.",
            "The defensible paper position is: Gemma 3 is project history; Gemma 4 is the quantitative study unless the Gemma 3 archive is recovered.",
        ],
        note="This keeps the deck honest and avoids inventing Gemma 3 numbers that are not preserved.",
    )

    add_bullets_slide(
        prs,
        "Experimental Timeline",
        [
            "V3b: early Gemma 4 flat extractor on product key JSON; it recovered many values but produced a noisy uncontrolled key space.",
            "V5: simpler flat JSON with IndiaMART-style keys plus cat74 coverage; still below qmeans on strict key+value quality.",
            "V6: stricter canonical supervision plus gold data; this is where Gemma 4 clearly beat qmeans on exact key+value match rate.",
            "V7 to V11: move to the nested priority schema; V7 became useful, V8 to V10 regressed, and V11 became the best standalone priority model.",
        ],
    )

    add_bullets_slide(
        prs,
        "How We Trained Gemma 4",
        [
            "V3b used product_train_with_keys data, max length 256, 3 epochs, and LoRA r=32 alpha=64 to learn early flat extraction.",
            "V5 simplified the target into flat JSON with IndiaMART-style keys and added cat74 motor-domain coverage.",
            "V6 tightened the problem again with a strict canonical prompt, cleaned training sets, gold_1k_v2 supervision, and remapped motor-domain keys.",
            "V11 used a balanced nested recipe: cat74 priority rows, V6 flat-to-nested rows, and gold_1k_v2 converted to nested form with a restrained learning rate of 7e-5.",
        ],
        note="Key insight: the best standalone result came from a balanced rollback, not from a larger or more aggressive recipe.",
    )

    add_table_slide(
        prs,
        "Early Gemma 4 V3b Evidence",
        ["System View", "Key F1", "Value Exact", "Takeaway"],
        [
            ["V3b raw (98-query view)", "13.7%", "57.2%", "Strong values, very noisy keys"],
            ["V3b normalized (98-query view)", "20.4%", "50.2%", "Better canonical control"],
            ["V3b normalized (full 1k view)", "20.5%", "42.5%", "Same pattern holds at scale"],
            ["qmeans (98-query view)", "45.0%", "19.4%", "Broad key coverage baseline"],
        ],
        note="Interpretation: the earliest Gemma 4 model already understood many values, but key discipline had to be learned in later stages.",
        col_widths=[3.0, 1.4, 1.6, 5.2],
    )

    add_table_slide(
        prs,
        "Flat Extraction Milestones",
        ["System", "Benchmark", "Headline Metric", "Value"],
        [
            ["qmeans_v2", "V5 compare", "KV strict F1", "37.23"],
            ["V5_flat", "V5 compare", "KV strict F1", "28.05"],
            ["qmeans", "Legacy 1k", "Exact KV match", "27.1%"],
            ["Gemma V6", "Legacy 1k", "Exact KV match", "40.7%"],
        ],
        note="Interpretation: Gemma 4 did not win immediately, but by the V6 stage it clearly surpassed qmeans on exact key+value match rate.",
        col_widths=[2.4, 2.3, 4.0, 1.8],
    )

    add_table_slide(
        prs,
        "Priority Benchmark Summary",
        ["Model", "Key F1", "Key+Value F1", "Main Role"],
        [
            ["qmeans", "45.50%", "6.00%", "Broad key coverage baseline"],
            ["V7", "29.09%", "15.52%", "First useful priority model"],
            ["V8", "15.88%", "10.72%", "Higher capacity, lower recall"],
            ["V9", "4.75%", "2.69%", "Heavy gold weighting collapse"],
            ["V10", "0.00%", "0.00%", "Clean-only collapse"],
            ["V11", "45.44%", "20.26%", "Best standalone priority model"],
            ["Hybrid V7+qmeans", "54.38%", "17.22%", "Best operational system"],
        ],
        note="V11 is the strongest standalone model. The hybrid remains the strongest deployment path for key F1.",
        col_widths=[2.8, 1.5, 1.8, 5.1],
    )

    add_bullets_slide(
        prs,
        "What V11 Changed",
        [
            "V11 key F1 reached 45.44%, effectively tying qmeans at 45.50% on broad key recovery.",
            "V11 key+value F1 reached 20.26%, beating qmeans at 6.00% and beating V7 at 15.52%.",
            "This shows the project recovered from the V8 to V10 regression and found a much stronger standalone model.",
            "The recovery came from a conservative data and adapter recipe, not from scaling up the model or oversampling gold more aggressively.",
        ],
        note="Short verdict: V11 is the best standalone model; hybrid is still the best overall system-level solution.",
    )

    add_example_slide(
        prs,
        "Concrete Example",
        "Query and Hybrid Output",
        "Query: siemens 45 hp 3 phase 1440 rpm motor\n\nHybrid final result:\nbrand=siemens\npart type=motor\npower=45 hp\nphase=3 phase\nrpm=1440 rpm\n\nNested output preserves attribute order and key-synonym order.",
        "Why The Example Matters",
        "The preserved V7 raw output for this query had extra nested wrappers under attribute_priorityN.\n\nAfter robust parsing, V7 and the hybrid recovered the correct flat semantics. qmeans recovered most values but weakened the part-type assignment in this case.\n\nThis example shows both the research challenge and the practical value of the system.",
    )

    add_bullets_slide(
        prs,
        "Problems We Solved",
        [
            "Moved the task from unordered extraction to an explicit priority-aware intent representation.",
            "Reduced open-key drift by pushing the model toward canonical attribute families instead of free-form key names.",
            "Proved Gemma 4 can beat qmeans on exact key+value quality in flat extraction by the V6 stage.",
            "Recovered a strong standalone priority model with V11 after several negative-result ablations.",
            "Built a hybrid inference path that combines structured Gemma output with qmeans coverage for the best operational outcome.",
        ],
    )

    add_bullets_slide(
        prs,
        "Why The Negative Results Matter",
        [
            "V8 showed that more capacity and more gold supervision did not automatically improve the structured model.",
            "V9 showed that heavier gold weighting could collapse recall even further.",
            "V10 showed that clean-only alignment could over-correct so strongly that extraction effectively disappeared.",
            "V11 showed that the right recovery path was better balance, not more force.",
            "That failure sequence strengthens the paper because it explains why the final recipe works.",
        ],
    )

    add_bullets_slide(
        prs,
        "What Is Still Needed For Publication",
        [
            "Add rank-aware metrics: attribute top-1 accuracy, pairwise order accuracy, key-synonym top-1 accuracy, and key-synonym MRR.",
            "Add structure metrics formally: parse rate, schema compliance rate, placeholder-key rate, and deep-wrapper rate.",
            "Add exact structured match as the strict downstream metric.",
            "Either recover the missing Gemma 3 benchmark archive or scope the final paper explicitly to the Gemma 4 program.",
        ],
        note="The paper story is already strong, but the evaluation stack still under-measures the full ranking contribution.",
    )

    add_bullets_slide(
        prs,
        "Team Takeaway",
        [
            "We now have a complete internal research story that starts with V3b, passes through the V6 flat win, and ends with the V11 standalone result plus the hybrid deployment result.",
            "The project has produced both positive findings and useful negative ablations, which strengthens the research argument rather than weakening it.",
            "The updated draft and deck are ready for internal team circulation with an honest Gemma 3 position and fuller experiment coverage.",
        ],
        note="Best standalone metric: V11 key+value F1 = 20.26%. Best overall operational key F1: Hybrid = 54.38%.",
    )

    prs.save(OUT)


if __name__ == "__main__":
    build()